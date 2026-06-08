"""
Build the MTech thesis-defense slide deck (.pptx).

Sparse bullet slides + full speaker-note scripts. Clean editable template.
Run:  python thesis/build_defense_deck.py
Output: thesis/thesis_defense.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------- theme
NAVY   = RGBColor(0x1F, 0x38, 0x64)   # titles / title-slide band
ACCENT = RGBColor(0x2E, 0x75, 0xB6)   # accent bar / highlights
DARK   = RGBColor(0x33, 0x33, 0x33)   # body text
GREY   = RGBColor(0x70, 0x70, 0x70)   # subtle text
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xEE, 0xF3, 0xF9)   # section-divider background
FONT   = "Calibri"

HERE  = os.path.dirname(os.path.abspath(__file__))
ROOT  = os.path.dirname(HERE)
SEVERITY_IMG = os.path.join(ROOT, "results", "task05", "_extracted_pngs", "cell3_out0.png")

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------- helpers
def _set(run, size, color=DARK, bold=False, italic=False, font=FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


def rect(slide, l, t, w, h, color, line=False):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if not line:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def title_bar(slide, title):
    """Accent stripe + slide title at top of a content slide."""
    rect(slide, Inches(0.0), Inches(0.0), SW, Inches(0.14), ACCENT)
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), SW - Inches(1.2), Inches(0.95))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    _set(r, 30, NAVY, bold=True)
    return tb


def content_slide(title, bullets, notes, img=None, img_caption=None):
    """
    bullets: list of (level, text). level 0 = main, 1 = sub, 2 = sub-sub.
    img: optional image path placed on the right / below.
    """
    s = prs.slides.add_slide(BLANK)
    title_bar(s, title)

    body_w = SW - Inches(1.2)
    if img:
        body_w = Inches(6.0)
    body = s.shapes.add_textbox(Inches(0.7), Inches(1.55), body_w, SH - Inches(2.1))
    tf = body.text_frame
    tf.word_wrap = True

    first = True
    for level, text in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        if level == 0:
            p.space_after = Pt(12)
            mark, size, color, bold = "▪  ", 21, DARK, False
        elif level == 1:
            p.space_after = Pt(7)
            mark, size, color, bold = "–  ", 18, RGBColor(0x44, 0x44, 0x44), False
        else:
            p.space_after = Pt(5)
            mark, size, color, bold = "·  ", 16, GREY, False
        p.space_before = Pt(2)
        r = p.add_run(); r.text = mark + text
        _set(r, size, color, bold=bold)

    if img and os.path.exists(img):
        s.shapes.add_picture(img, Inches(6.7), Inches(1.7), width=Inches(6.2))
        if img_caption:
            cap = s.shapes.add_textbox(Inches(6.7), Inches(5.3), Inches(6.2), Inches(0.6))
            cp = cap.text_frame; cp.word_wrap = True
            r = cp.paragraphs[0].add_run(); r.text = img_caption
            _set(r, 12, GREY, italic=True)
    add_notes(s, notes)
    return s


def big_image_slide(title, img, notes, caption=None):
    s = prs.slides.add_slide(BLANK)
    title_bar(s, title)
    if os.path.exists(img):
        s.shapes.add_picture(img, Inches(0.7), Inches(1.7), width=Inches(11.9))
    if caption:
        cap = s.shapes.add_textbox(Inches(0.7), Inches(6.4), Inches(11.9), Inches(0.6))
        r = cap.text_frame.paragraphs[0].add_run(); r.text = caption
        _set(r, 13, GREY, italic=True)
    add_notes(s, notes)
    return s


def table_slide(title, headers, rows, notes, bold_last_row=True, col_widths=None, footnote=None):
    s = prs.slides.add_slide(BLANK)
    title_bar(s, title)
    nrows, ncols = len(rows) + 1, len(headers)
    tw = SW - Inches(1.4)
    th = Inches(0.55) + Inches(0.42) * len(rows)
    gfx = s.shapes.add_table(nrows, ncols, Inches(0.7), Inches(1.7), tw, th)
    tbl = gfx.table
    if col_widths:
        for i, frac in enumerate(col_widths):
            tbl.columns[i].width = int(tw * frac)
    # header
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
        r = p.add_run(); r.text = h
        _set(r, 15, WHITE, bold=True)
    # body
    for ri, row in enumerate(rows, start=1):
        last = (ri == len(rows))
        for c, val in enumerate(row):
            cell = tbl.cell(ri, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT if ri % 2 else WHITE
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            r = p.add_run(); r.text = val
            _set(r, 14, NAVY if (last and bold_last_row) else DARK,
                 bold=(last and bold_last_row))
    if footnote:
        fn = s.shapes.add_textbox(Inches(0.7), Inches(1.75) + th, SW - Inches(1.4), Inches(0.6))
        r = fn.text_frame.paragraphs[0].add_run(); r.text = footnote
        _set(r, 13, GREY, italic=True)
    add_notes(s, notes)
    return s


def section_slide(kicker, title, notes):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, LIGHT)
    rect(s, Inches(0.0), Inches(3.05), SW, Inches(0.10), ACCENT)
    kb = s.shapes.add_textbox(Inches(1.0), Inches(2.2), SW - Inches(2.0), Inches(0.6))
    r = kb.text_frame.paragraphs[0].add_run(); r.text = kicker.upper()
    _set(r, 16, ACCENT, bold=True)
    tb = s.shapes.add_textbox(Inches(1.0), Inches(3.3), SW - Inches(2.0), Inches(1.6))
    tf = tb.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = title
    _set(r, 40, NAVY, bold=True)
    add_notes(s, notes)
    return s


# ================================================================ SLIDES

# 1 -------------------------------------------------- TITLE
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, WHITE)
rect(s, Inches(0.0), Inches(0.0), Inches(0.35), SH, NAVY)
rect(s, Inches(0.0), Inches(4.55), SW, Inches(0.06), ACCENT)
tb = s.shapes.add_textbox(Inches(1.0), Inches(1.7), SW - Inches(1.8), Inches(2.6))
tf = tb.text_frame; tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = "Refining the SmoothQuant Calibration Step for Efficient INT8 Inference of Large Language Models"
_set(r, 38, NAVY, bold=True)
sub = s.shapes.add_textbox(Inches(1.0), Inches(4.75), SW - Inches(1.8), Inches(1.0))
tf = sub.text_frame; tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = "Quantile-based smoothing and per-layer migration strength for post-training W8A8 quantization"
_set(r, 19, GREY, italic=True)
foot = s.shapes.add_textbox(Inches(1.0), Inches(6.05), SW - Inches(1.8), Inches(1.2))
tf = foot.text_frame; tf.word_wrap = True
for i, line in enumerate([
    "Anand  ·  MTech Thesis Defense",
    "Department of [Department], IIT Guwahati  ·  2026",
    "Advisor: [Advisor Name]",
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run(); r.text = line
    _set(r, 16, DARK, bold=(i == 0))
add_notes(s, """
Good morning. Thank you for being here.
My thesis is on post-training quantization of large language models.
The goal is cheaper inference. Run these models on smaller hardware without losing accuracy.
I build on SmoothQuant. I look closely at one stage of it, the calibration step, and I propose two principled changes there.
Over the next half hour I will walk through the problem, the method I build on, the three refinements I propose, and the results across five OPT sizes and two other architectures.
[Fill in your department and advisor name on this slide.]
""")

# 2 -------------------------------------------------- OUTLINE
content_slide(
    "Outline",
    [(0, "The problem: models have outgrown the hardware"),
     (0, "Quantization, and why W8A8 is hard"),
     (0, "SmoothQuant: the method I build on"),
     (0, "Three refinements to the calibration step"),
     (0, "Results: five OPT sizes, two more architectures, real INT8 memory"),
     (0, "Conclusion and future work")],
    """
Here is the plan for the talk.
First I motivate why inference cost is the real bottleneck today.
Then quantization, and the specific reason eight-bit activations are hard.
Then SmoothQuant, which is the method I start from.
The core of the talk is three refinements I make to its calibration step.
Then the results. I evaluate on five OPT sizes, then on Falcon and Llama-2, and finally I measure real memory on hardware.
I close with conclusions and where this can go next.
""")

# 3 -------------------------------------------------- SECTION: motivation
section_slide("Part 1", "Motivation: the cost of inference",
"""
Let me start with why this problem matters.
""")

# 4 -------------------------------------------------- problem
content_slide(
    "Models have outgrown the hardware",
    [(0, "Parameters: 100M → 100B+ in five years"),
     (0, "GPU memory has barely doubled in the same time"),
     (0, "A 7B model in FP16 already exceeds most consumer GPUs"),
     (0, "Inference is the recurring cost"),
     (1, "Training is paid once"),
     (1, "Inference is paid on every single input"),
     (0, "Bottleneck for laptops, phones, edge devices")],
    """
The scale problem is well known. Model size has grown by roughly a thousand times in five years. GPU memory has not kept up. It has maybe doubled.
So even a seven billion parameter model in full precision does not fit on a typical consumer GPU.
The key point is this. Training is a one-time cost. Inference is paid every time a user sends an input.
So if we want these models to run outside a data centre, on a laptop or a phone, the thing we have to reduce is the cost of inference.
That is the setting of this thesis.
""")

# 5 -------------------------------------------------- quantization
content_slide(
    "Quantization is the lever",
    [(0, "Replace FP16 tensors with 8-bit integers"),
     (1, "~half the memory, ~2× the throughput"),
     (1, "Maps to real hardware: INT8 Tensor Cores, AMX, DSPs"),
     (0, "Post-training (PTQ), not retraining (QAT)"),
     (1, "One calibration pass on a small sample. Training-free."),
     (0, "W8A8: both weights and activations in INT8"),
     (1, "INT8 GEMM kernels need both operands in INT8"),
     (1, "Weight-only saving shrinks as the model grows")],
    """
Quantization is one of the few techniques that attacks memory, bandwidth, and compute at the same time. We replace sixteen-bit floats with eight-bit integers. Roughly half the memory and twice the throughput. And it maps onto hardware that already exists, the integer tensor cores on modern GPUs.
There are two ways to do it. Quantization-aware training retrains the model. At this scale that is as expensive as pre-training, so it is not practical. I use post-training quantization. One calibration pass on a small sample, no retraining.
Now an important detail. Early in the thesis I measured this myself. If you only quantize the weights, you do not get the fifty percent memory saving you expect. And it gets worse as the model grows. The reason is that activations, the intermediate tensors, come to dominate memory. So we must quantize both. Weights and activations, both in eight bits. That is W8A8. It is also what the integer kernels require, both operands must be integer.
""")

# 6 -------------------------------------------------- outliers
content_slide(
    "The catch: activation outliers",
    [(0, "Weights are flat. Activations are not."),
     (0, "Outliers are structured (Dettmers et al., LLM.int8())"),
     (1, "Sparse: ~0.1% of channels"),
     (1, "Persistent: an outlier channel is large for every token"),
     (0, "Per-tensor activation scale is set by the largest outlier"),
     (1, "Typical channels collapse to 2–3 usable integer levels"),
     (0, "Naive W8A8 collapses: OPT-175B accuracy 71.6% → 32.3%")],
    """
Here is why eight-bit activations are hard.
Weights have a flat, well-behaved distribution. Activations do not. They carry extreme outliers.
Dettmers and colleagues studied these carefully. Two facts matter. The outliers are sparse, about a tenth of a percent of channels. And they are persistent. If a channel is an outlier channel, it is large for every token, not just a few.
Now think about a single per-tensor scale. That scale is dictated by the largest outlier in the whole tensor. Every normal channel is then squeezed into the small range that is left. Sometimes only two or three integer levels remain for a typical channel. That destroys the information.
The result is catastrophic. On the largest OPT model, naive eight-bit drops zero-shot accuracy from seventy-two percent down to thirty-two. That is close to random guessing.
So the outlier problem is the whole problem.
""")

# 7 -------------------------------------------------- smoothquant
content_slide(
    "SmoothQuant: migrate the difficulty",
    [(0, "Per-channel activation quant would fix it, but breaks INT8 kernels"),
     (0, "SmoothQuant's idea: move the difficulty into the weights"),
     (1, "Y = X·W = (X·diag(s)⁻¹)·(diag(s)·W)  — same output"),
     (1, "Flatten activation channels. Weights absorb the variance."),
     (0, "Scale s fixed once at calibration, fused into LayerNorm"),
     (1, "No extra kernel at inference"),
     (0, "One knob: migration strength α ∈ [0,1]"),
     (0, "Strongest training-free, all-INT8 baseline today")],
    """
SmoothQuant is the method I build on. Here is the idea.
The statistically correct fix would be per-channel activation quantization, a separate scale per channel. But that cannot be fused into an integer matrix multiply, so it is not practical on hardware.
SmoothQuant takes a different route. It uses a simple algebraic fact. You can divide the activation by a per-channel scale s and multiply the weight by the same s, and the output of the layer does not change. So you choose s to flatten the activation channels. The difficulty does not disappear. It moves onto the weights. And that is a good trade, because weights start out flat and have room to absorb it.
The scale s is computed once during calibration and folded into the preceding LayerNorm. So nothing extra runs at inference time.
There is exactly one knob. The migration strength alpha, between zero and one. It controls how much difficulty you move. Alpha equal to one half is the value the paper uses for OPT.
This is the strongest training-free, fully integer baseline available. That is why I start here.
""")

# 8 -------------------------------------------------- SECTION: proposed
section_slide("Part 2", "Where the baseline is coarse — and three refinements",
"""
SmoothQuant works. But when I looked closely at the calibration step, I found three places where it makes a coarse choice. That is the contribution of this thesis.
""")

# 9 -------------------------------------------------- the gap
content_slide(
    "Three coarse choices in the baseline",
    [(0, "The baseline treats the network as statistically uniform"),
     (0, "1.  Weights quantized per-tensor"),
     (1, "but smoothing pushes variance onto the weights"),
     (0, "2.  Scale uses the per-channel max"),
     (1, "one anomalous token can set the scale forever"),
     (0, "3.  One global α for the whole network"),
     (1, "and re-grid-searched for every new model"),
     (0, "The real outlier distribution is not uniform. Three fixes follow.")],
    """
When I studied the calibration step, I found three assumptions that are coarser than the data warrants. In each case the baseline treats the network as if it were uniform, and it is not.
First, the weights are quantized with a single per-tensor scale. But smoothing deliberately pushes variance onto the weights. So that is exactly where a single scale hurts.
Second, the scale comes from the maximum of each channel. The maximum is the single most extreme token. So one freak token in the calibration data sets the scale for that channel for the entire deployment.
Third, alpha is one global number for the whole network. And the paper re-tunes it per model with a grid search.
Each of these is a mismatch with the real distribution. I propose one refinement for each. Importantly, all three live entirely in the offline calibration step. Nothing changes at inference.
""")

# 10 ------------------------------------------------- refinement 1
content_slide(
    "Refinement 1 — per-channel weight quantization",
    [(0, "Smoothing moves variance into the weights"),
     (0, "A per-tensor weight scale then squeezes typical channels"),
     (1, "the same failure, now on the weight side"),
     (0, "Fix: one scale per output channel of the weight"),
     (0, "Still compatible with INT8 GEMM kernels"),
     (1, "output-channel is an outer dimension of the matmul"),
     (1, "rescaling is a post-GEMM vector multiply"),
     (0, "Paper uses it silently on Llama/Falcon — never named or compared on OPT")],
    """
The first refinement. Per-channel weight quantization. I call this configuration C in the experiments.
The logic follows directly from how smoothing works. Smoothing moves variance onto the weights. If you then quantize the weights with one global scale, the typical weight channels get squeezed, exactly the same failure we saw on the activation side. So the natural countermeasure is to give each output channel of the weight its own scale.
And critically, this is free on hardware. The output-channel dimension is an outer dimension of the matrix multiply. So the per-channel scale is just a vector multiply applied after the integer GEMM. It never enters the integer accumulation. The kernels are unchanged.
Here is the interesting part. The SmoothQuant paper actually does use per-channel weights, but only for its later Llama and Falcon numbers, and it mentions it in one passing sentence. It never names it, and it never compares it against the per-tensor schemes on OPT. I make that comparison explicit.
""")

# 11 ------------------------------------------------- refinement 2
content_slide(
    "Refinement 2 — quantile-based smoothing",
    [(0, "Outliers also occur within a channel, across tokens"),
     (0, "Baseline uses max — hypersensitive to one spike"),
     (0, "Replace max with a high quantile  Q_p  (p ≈ 0.999)"),
     (1, "ignores the single-token spike, keeps the true scale"),
     (1, "p = 1 recovers the original formula exactly"),
     (0, "Same problem the paper patched ad hoc (top-2% clip on GLM-130B)"),
     (1, "now solved at the smoothing step, in the math"),
     (0, "Computed exactly via a per-channel top-K buffer")],
    """
The second refinement is at the inner calibration step. Quantile-based smoothing.
SmoothQuant solves the problem of outliers across channels. But within a single channel, across tokens, the distribution is also not uniform. Most tokens sit near the median, but a few spike well above it. And the baseline uses the maximum. So one anomalous token in the calibration sample sets the scale for that channel for every forward pass afterwards.
My fix is to replace the maximum with a high percentile. The 99.9th percentile, for example. It ignores the very top spike and reports the true typical scale of the channel. And note, if you set the percentile to one, you get the original formula back exactly. So this is a strict generalisation.
Here is the nice part. The SmoothQuant authors hit this exact problem on GLM-130B. They patched it by clipping the top two percent of tokens. They themselves call it a workaround, and it lives at the quantization step. I solve the same problem properly, in the smoothing formula itself, with no separate clip to tune.
One engineering note. Storing every activation to compute a percentile is infeasible, hundreds of gigabytes. I keep an exact per-channel top-K buffer instead. The result is identical to full storage, at a fraction of the memory.
""")

# 12 ------------------------------------------------- refinement 3
content_slide(
    "Refinement 3 — per-layer migration strength α(l)",
    [(0, "Outlier severity is not constant across depth"),
     (1, "severity σ(l) = max-channel / median-channel, per layer"),
     (0, "Profile is sharply structured, same shape across sizes"),
     (1, "layer 0 low · layers 1–3 peak · smooth decay after"),
     (0, "One global α under-smooths hard layers, over-smooths easy ones"),
     (0, "Map severity → α(l):   α(l) = α_min + (α_max − α_min)·σ̃(l)"),
     (1, "two scalars replace the per-model grid search"),
     (1, "offline only — zero inference cost")],
    """
The third refinement is at the network level. A per-layer alpha.
The migration strength is the central knob, and the baseline uses one value for the whole network. But the severity of outliers is not constant across depth. I measure it directly. For each layer, severity is the ratio of the worst channel to the median channel.
And the profile is sharply structured. I will show you the figure on the next slide. Layer zero is low. Layers one to three hold a peak. Then it decays smoothly. The same shape repeats at every model size.
So a single global alpha is wrong. It under-smooths the hard layers near the peak and over-smooths the easy layers in the tail.
My fix is a simple linear map. Take the severity, normalise it within the model, and map it onto an alpha between a lower and an upper bound. The hardest layer gets the highest alpha, the easiest gets the lowest, everything else interpolates.
The key point is what this replaces. The paper runs a grid search per model to find one alpha. I replace that whole search with two fixed scalars and a profile I already measured during calibration. No search, and again, purely offline.
""")

# 13 ------------------------------------------------- THE CLAIM
content_slide(
    "What the contribution actually is",
    [(0, "Not “we beat SmoothQuant” — the PPL deltas are within noise"),
     (0, "The claim is structural:"),
     (1, "α(l) matches the paper's per-model tuned α — with no grid search"),
     (0, "Paper Table 7 tunes α per model: 0.6 → 0.9"),
     (1, "Falcon-7B 0.6 · Mistral 0.8 · Llama-2-7B/13B 0.85 · Llama-70B 0.9"),
     (1, "each row implies a separate search"),
     (0, "Same construction across architectures — saves compute and researcher time"),
     (0, "All three refinements are offline → zero deployment overhead")],
    """
This is the most important slide. I want to be precise and honest about what I am claiming.
I am not claiming I beat SmoothQuant on accuracy. The perplexity differences are real but small, within evaluation noise. If I stood here and said my method gives lower perplexity, that would be overclaiming.
The claim is structural. My per-layer alpha matches the paper's carefully tuned per-model alpha, without doing any tuning.
Look at Table 7 of the SmoothQuant paper. They report a different hand-tuned alpha for every model. Zero point six for Falcon. Zero point eight five for Llama-2. Zero point nine for the seventy-billion model. Every one of those numbers implies a separate grid search on that model.
My method derives alpha from the severity profile I already measured during calibration. The same construction, with the same two scalars, across every architecture. No per-model search.
So the contribution is that I reach the same quality the paper reaches, but I remove the per-model tuning cost. The perplexity numbers are the evidence. The removal of the search is the claim. And because everything is offline, there is zero cost at deployment.
""")

# 14 ------------------------------------------------- SECTION: results
section_slide("Part 3", "Results",
"""
Now the evidence. Let me show what these refinements actually do across models.
""")

# 15 ------------------------------------------------- setup
content_slide(
    "Experimental setup",
    [(0, "Models: OPT 125M · 1.3B · 2.7B · 6.7B · 13B (+ Falcon-7B, Llama-2-7B)"),
     (0, "Calibration: 512 Pile sentences × 512 tokens"),
     (0, "Accuracy: WikiText-2 perplexity @ seq-len 2048"),
     (1, "+ 7 zero-shot tasks via lm-evaluation-harness"),
     (0, "Two regimes:"),
     (1, "fake-quant for all accuracy numbers"),
     (1, "real INT8 (torch-int CUTLASS) for the memory measurement"),
     (0, "Hardware: Google Colab T4 / L4 / A100")],
    """
The setup, briefly.
I focus on the OPT family, five sizes from 125 million to 13 billion. That is the same family SmoothQuant uses for its main scaling results, so it is directly comparable. I then add Falcon and Llama-2 to test other architectures.
Calibration is a fixed sample of the Pile, 512 sentences of 512 tokens.
The main accuracy metric is WikiText-2 perplexity at sequence length 2048. That sequence length matters. It is the convention in both SmoothQuant and LLM.int8. I never shorten it to fit memory, because that would break comparability. I also run seven zero-shot tasks to confirm the perplexity gains carry to real tasks.
Two regimes. For accuracy I use fake quantization, the standard way to evaluate a method. For the memory claim, I run the real integer kernels on hardware.
Everything runs on Colab GPUs, the T4, L4, and A100.
""")

# 16 ------------------------------------------------- result quantile PPL
table_slide(
    "Result 1 — quantile smoothing recovers FP16",
    ["OPT model", "FP16", "Naive W8A8", "Max O1 (α=0.5)", "Ours: quantile + PCW"],
    [["125M",  "27.57", "30.23", "28.31", "27.63"],
     ["1.3B",  "14.47", "15.59", "14.83", "14.62"],
     ["2.7B",  "12.34", "13.47", "12.39", "12.34"]],
    """
The first result. This is quantile smoothing combined with per-channel weights, against the baselines, on three OPT sizes.
Read any row. The FP16 column is the ceiling, what we want to reach. Naive eight-bit is the floor, what happens with no smoothing. The O1 column is the SmoothQuant baseline at its tuned alpha.
Our column, on the right, is essentially on top of FP16 in every row. On 125M, 27.63 against an FP16 of 27.57, near lossless. On 2.7B, 12.34 against 12.34, the difference is in the fourth decimal.
And in every row we are below the O1 baseline. The gains are small in absolute terms because these mid-size models have headroom, but the trend is consistent across scale. That consistency is the point. It is not a fluke on one model.
""",
    col_widths=[0.18, 0.18, 0.20, 0.22, 0.22],
    footnote="WikiText-2 perplexity, seq-len 2048. Lower is better. Ours: Q_p with p≈0.999, per-channel weights.")

# 17 ------------------------------------------------- zero-shot
table_slide(
    "Result 1b — gains carry to downstream tasks",
    ["OPT-1.3B config", "LAMBADA", "HellaSwag", "PIQA", "WinoG.", "RTE", "Avg."],
    [["FP16",            "57.87", "53.70", "71.71", "59.43", "52.35", "58.32"],
     ["Naive W8A8",      "54.32", "51.74", "70.24", "57.14", "56.68", "57.25"],
     ["Max O1 (α=0.5)",  "57.54", "53.42", "71.16", "58.80", "51.99", "58.10"],
     ["Ours: quantile",  "57.44", "53.65", "71.16", "59.27", "52.71", "58.29"]],
    """
A perplexity win is only meaningful if it shows up on real tasks. So I ran seven zero-shot benchmarks on OPT-1.3B. I am showing five of the columns here plus the average.
The bottom row is our method. The average is 58.29 percent. The FP16 ceiling is 58.32. So we are within three hundredths of a point of full precision. We are above the O1 baseline at 58.10, and well above naive eight-bit at 57.25.
So the perplexity story and the task-accuracy story agree. The method does not just improve a language-modelling number, it preserves downstream quality.
""",
    col_widths=[0.24, 0.135, 0.145, 0.115, 0.125, 0.11, 0.13],
    footnote="Zero-shot accuracy (%), lm-evaluation-harness. Full seven-task table in the thesis. Higher is better.")

# 18 ------------------------------------------------- severity figure
big_image_slide(
    "Result 2 — the per-layer severity profile",
    SEVERITY_IMG,
    """
This is the figure behind the third refinement, and it is the one to spend a moment on.
On the x-axis is the layer index, from the first decoder layer to the last. On the y-axis, on a log scale, is the severity. The ratio of the worst channel to the median channel at that layer. The two lines are the two sites I hook, the attention input and the feed-forward input.
Three things to notice. First, the shape. Layer zero is low. Layers one to three jump to a sharp peak. Then it decays smoothly to a tail. Second, this same shape appears in all three model sizes, side by side. It is not noise, it is a stable property of the architecture. Third, the two lines, attention and feed-forward, lie almost exactly on top of each other.
And the magnitude grows with scale. The peak severity is about thirteen times at 2.7B, twenty-three times at 6.7B, and forty-six times at 13B. So the bigger the model, the more wrong a single global alpha becomes. This is precisely why a per-layer alpha helps more as we scale up. The structure is here, sitting in the calibration data, free to use.
""",
    caption="Per-layer outlier severity across OPT-2.7B / 6.7B / 13B. Same shape at every scale; peak grows 13× → 23× → 46×.")

# 19 ------------------------------------------------- per-layer alpha PPL
table_slide(
    "Result 2b — per-layer α vs global α",
    ["OPT model", "FP16", "Global α (O1, 0.5)", "Ours: per-layer α(l)"],
    [["125M", "27.57", "28.31", "27.59"],
     ["1.3B", "14.47", "14.83", "14.63"],
     ["2.7B", "12.34", "12.39", "12.37"]],
    """
Here is the per-layer alpha on its own, against the standard fixed global alpha.
Again read any row. Our per-layer alpha, on the right, sits just above FP16 and below the global-alpha baseline in every case. On 125M, 27.59 against an FP16 of 27.57. Almost exactly lossless.
Remember the framing. I am not claiming a large win here. I am claiming I match the tuned baseline without tuning anything. These numbers come from two fixed scalars and the severity profile. There was no grid search on any of these models.
""",
    col_widths=[0.22, 0.22, 0.28, 0.28],
    footnote="WikiText-2 PPL, seq-len 2048. α(l) from severity profile, α_min=0.5, α_max=0.9. No per-model tuning.")

# 20 ------------------------------------------------- scale robustness
table_slide(
    "Result 3 — robustness at the outlier-heavy scale",
    ["OPT model", "FP16", "Naive W8A8", "Max O1", "Ours: per-layer α(l)"],
    [["6.7B", "10.67", "25.91", "10.70", "10.68"],
     ["13B",  "9.94",  "4325.68", "10.18", "9.98"]],
    """
Now the most important test. LLM.int8 showed that activation outliers become severe above roughly six billion parameters. That is the regime where naive eight-bit collapses. So this is where a quantization method has to prove itself.
Look at the naive W8A8 column. At 6.7B perplexity jumps to twenty-six, more than double FP16. At 13B it explodes to over four thousand. The model is effectively destroyed. That single number is the clearest picture of how serious the outlier problem is at scale.
Now look at our column. At 6.7B, 10.68 against an FP16 of 10.67. At 13B, 9.98 against 9.94. We hold right on the FP16 ceiling, and we stay below the O1 baseline, exactly in the regime where naive quantization falls apart.
So the method does not just work on small models. It works precisely where it has to.
""",
    col_widths=[0.18, 0.16, 0.24, 0.18, 0.24],
    footnote="WikiText-2 PPL, seq-len 2048. The 13B naive collapse (4325) is the outlier problem in one number.")

# 21 ------------------------------------------------- cross-arch
table_slide(
    "Result 4 — same construction, other architectures",
    ["Model", "FP16", "Max + paper-tuned α", "Ours: per-layer α(l)"],
    [["Falcon-7B  (paper α=0.60)", "6.95", "6.98", "6.99"],
     ["Llama-2-7B (paper α=0.85)", "5.82", "5.87", "5.86"]],
    """
This slide is where the central claim is tested. So far everything was OPT. Falcon and Llama-2 are different architectures, and crucially the paper uses very different tuned alphas for them. Zero point six for Falcon. Zero point eight five for Llama-2. Those came from per-model searches.
The middle column is the paper's tuned baseline, using its hand-picked alpha for that model. The right column is my per-layer alpha, with no tuning, the same construction I used on OPT.
On Falcon, 6.99 against the tuned 6.98. On Llama-2, 5.86 against the tuned 5.87. They match to within a hundredth of a perplexity point.
This is the proof of the claim. My method reaches the same quality as the paper's per-model tuned alpha, on two architectures with very different alpha needs, without running a single grid search. That is the contribution, demonstrated.
""",
    col_widths=[0.34, 0.16, 0.26, 0.24],
    footnote="WikiText-2 PPL, seq-len 2048. α-bounds set once per family from the severity spread, not tuned on the eval set.")

# 22 ------------------------------------------------- memory
table_slide(
    "Result 5 — real INT8: free at deployment",
    ["OPT-1.3B", "Weight size", "Peak VRAM", "WikiText-2 PPL"],
    [["FP16",                       "2510 MB", "4557 MB", "14.62"],
     ["INT8 paper (max, α=0.5)",    "1358 MB", "3141 MB", "18.02"],
     ["INT8 ours (Q0.999, α=0.5)",  "1357 MB", "3140 MB", "17.96"]],
    """
The last result moves from fake quantization to the real integer kernels on an A100. This checks two things. That the memory saving is real, and that my change costs nothing at deployment.
Compared to FP16, the weights drop to fifty-four percent and peak GPU memory drops to sixty-nine percent. So the integer deployment delivers the saving, as expected.
Now compare the two integer rows, the paper's and mine. They are identical on memory to within four hundredths of a percent. Same model shape, same kernels, same footprint. And on perplexity my row is slightly better, 17.96 against 18.02.
So the conclusion is clean. My refinement is a pure calibration change. It produces the same integer model, the same memory, the same kernels, and it gives equal or better accuracy. It is free at deployment.
One honest note. The perplexity here, around eighteen, is worse than the fourteen-point-six we saw in fake-quant. That is because these particular kernels only support the coarsest per-tensor static scheme. The gap is the kernels, not the method. I will come back to that in future work.
""",
    col_widths=[0.34, 0.22, 0.22, 0.22],
    footnote="torch-int CUTLASS, O3 scheme. Both INT8 rows identical on memory to within 0.04%.")

# 23 ------------------------------------------------- challenges (defense-friendly)
content_slide(
    "Engineering challenges (selected)",
    [(0, "Building torch-int on a 2026 Colab took five separate patches"),
     (1, "a 2023-era CUDA extension on current PyTorch / CUDA 12"),
     (0, "The Pile mirror the repo points to is dead"),
     (1, "rebuilt the calibration set from a HuggingFace mirror"),
     (0, "datasets ≥ 3.0 silently broke half the zero-shot tasks"),
     (0, "T4 cannot host OPT-6.7B at seq-len 2048 — had to move to A100"),
     (0, "Llama-2-13B mirror ships a broken tokenizer (PPL ≈ vocab size)")],
    """
I want to spend one slide on the engineering, because a lot of the real work was here and it is fair game in a defense.
The integer kernels I needed, torch-int, were last updated in 2023. Getting them to build on a current Colab took five separate patches, from submodule URLs to C++ standard flags. Budget half a day, I learned.
The dataset the repository points to for calibration is offline, the host shut down years ago. I rebuilt the calibration set from a HuggingFace mirror.
A version bump in the datasets library silently broke several of the zero-shot tasks, so I had to pin it for reproducibility.
The free T4 GPU cannot fit the 6.7B model at the required sequence length, which is why those runs are on the A100. I refused to shorten the sequence length, because that would make the numbers incomparable.
And my favourite. The Llama-2-13B mirror ships a broken tokenizer. The model loaded fine and even produced text, but perplexity came out around the vocabulary size, the signature of random guessing. The lesson, when perplexity equals vocab size, suspect the tokenizer before the model. I fixed it by borrowing the 7B tokenizer, which is bit-identical.
These are all in the thesis. They are the reproducibility story behind the numbers.
""")

# 24 ------------------------------------------------- SECTION conclusion
section_slide("Part 4", "Conclusion", "")

# 25 ------------------------------------------------- conclusion
content_slide(
    "Conclusion",
    [(0, "Two principled fixes to the SmoothQuant calibration step"),
     (1, "quantile smoothing — handles within-channel outliers"),
     (1, "per-layer α(l) — handles severity drift across depth"),
     (1, "paired with per-channel weight quantization"),
     (0, "Training-free · INT8-compatible · zero inference cost"),
     (0, "Recovers FP16 quality across OPT 125M–13B, Falcon-7B, Llama-2-7B"),
     (0, "Matches paper-tuned quality with no per-model α search"),
     (0, "Narrow but clean: not a new architecture, a reproducible calibration fix")],
    """
To conclude.
I made two principled changes to one stage of SmoothQuant, the calibration step. Quantile smoothing, which handles outliers within a channel. And a per-layer alpha, which handles the way outlier severity drifts across depth. Both are paired with per-channel weight quantization, which follows naturally from how smoothing redistributes variance.
All of it is training-free, fully integer-compatible, and adds nothing at inference. The cost is paid once, offline.
Empirically, it recovers full-precision quality across five OPT sizes and on two other architectures. And it does so without the per-model alpha search the original method relies on.
I want to be measured about the size of the claim. This is not a new architecture and it is not a new training method. It is a narrow, clean, reproducible fix for two places where the SmoothQuant calibration leaves accuracy on the table. That is exactly what it sets out to be.
""")

# 26 ------------------------------------------------- future work
content_slide(
    "Future work",
    [(0, "W4A8 — quantile smoothing is bit-width independent"),
     (1, "combine with INT4 weight quantizers (GPTQ, AWQ)"),
     (0, "KV-cache quantization"),
     (1, "the last large FP16 consumer in the inference path"),
     (0, "Scale beyond 30B — needs tensor-parallel evaluation"),
     (0, "A per-channel-weight INT8 kernel"),
     (1, "would close the fake-quant → real-INT8 accuracy gap")],
    """
A few clear directions follow from this.
The first is four-bit weights with eight-bit activations. The smoothing step does not care about the bit width, so it should combine cleanly with four-bit weight quantizers like GPTQ and AWQ.
The second is the KV-cache. After weights and activations, the cache is the last large full-precision consumer of memory, and the same offline calibration philosophy applies.
The third is scale. Going above thirty billion parameters needs a multi-GPU, tensor-parallel evaluation setup. Verifying the severity profile holds there is the natural next step.
And the fourth ties back to the memory slide. The accuracy gap I saw in real integer mode is because the available kernels only support the coarse per-tensor scheme. Writing a proper per-channel-weight integer kernel would close the gap between the fake-quant accuracy I report and the real deployed accuracy.
""")

# 27 ------------------------------------------------- thank you
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, NAVY)
rect(s, Inches(0.0), Inches(4.35), SW, Inches(0.06), ACCENT)
tb = s.shapes.add_textbox(Inches(1.0), Inches(2.7), SW - Inches(2.0), Inches(1.6))
r = tb.text_frame.paragraphs[0].add_run(); r.text = "Thank you"
_set(r, 48, WHITE, bold=True)
sub = s.shapes.add_textbox(Inches(1.0), Inches(4.6), SW - Inches(2.0), Inches(1.0))
r = sub.text_frame.paragraphs[0].add_run(); r.text = "Questions and discussion"
_set(r, 22, RGBColor(0xCF, 0xDD, 0xEE))
add_notes(s, """
Thank you. I am happy to take questions.
Backup slides follow for the common ones: why not per-channel activations, why alpha changes between fake-quant and real INT8, and how the exact quantile buffer works.
""")

# ---- BACKUP --------------------------------------------------------------
section_slide("Backup", "Backup slides", "Backup material for likely questions.")

# B1
content_slide(
    "Backup — why not per-channel activation quantization?",
    [(0, "Statistically it is the ideal fix for activation outliers"),
     (0, "But the scale would sit on the inner (contraction) dimension"),
     (1, "it cannot be pulled out of the integer accumulation"),
     (1, "breaks the single fused INT8 GEMM"),
     (0, "SmoothQuant's whole point: move difficulty to a place hardware allows"),
     (0, "Per-token activation + per-channel weight both live on outer dims")],
    """
A common question. If outliers are per channel, why not just quantize activations per channel?
Statistically that is the ideal answer. The problem is purely about hardware. A per-channel activation scale sits on the inner dimension of the matrix multiply, the dimension that gets summed over. You cannot pull a scale out of that summation, so it has to live inside the integer accumulation, and that breaks the single fused integer kernel.
That is exactly the constraint SmoothQuant is built around. It moves the difficulty onto the weights, where the scale sits on an outer dimension and can be applied after the GEMM. Per-token activation scales and per-channel weight scales both have that property. That is why they are allowed and per-channel activation is not.
""")

# B2
content_slide(
    "Backup — α changes between fake-quant and real INT8",
    [(0, "Fake-quant winner on OPT-1.3B: α ≈ 0.9"),
     (0, "Real-INT8 (torch-int O3) winner: α = 0.5"),
     (0, "Why: O3 forces per-tensor weights + per-tensor static activations"),
     (1, "aggressive smoothing punishes the single per-tensor weight scale"),
     (0, "The optimum is regime-dependent — it must be re-checked per regime"),
     (0, "Quantile choice p stays fixed; only α shifts")],
    """
Another likely question, about a subtlety in the deployment numbers.
In fake-quant, with per-channel weights, the best alpha on 1.3B is around zero point nine. High smoothing is fine because per-channel weights absorb it.
But in the real integer deployment, the kernels force per-tensor weights and per-tensor static activations. There, the best alpha drops to zero point five. The reason is that with a single per-tensor weight scale, aggressive smoothing piles variance onto the weights and there is no per-channel scale to absorb it, so it hurts.
The lesson is that the optimal alpha is regime-dependent. The fake-quant optimum does not automatically transfer to the deployment regime. The percentile choice does transfer, only alpha shifts. This is documented honestly in the thesis rather than hidden.
""")

# B3
content_slide(
    "Backup — the exact-quantile top-K buffer",
    [(0, "Naive quantile = store every activation → hundreds of GB"),
     (0, "Instead: keep the top-K largest |values| per channel, sorted"),
     (1, "K = ⌈(1 − p_min)·N_total⌉,  p_min = 0.90"),
     (0, "Updated each forward pass via concat + torch.topk"),
     (0, "Read Q_p at rank ⌈p·(N−1)⌉ — exact, not approximate"),
     (1, "one pass yields p ∈ {0.90, 0.95, 0.99, 0.995, 0.999}"),
     (0, "Buffer cost: ~5 GB (1.3B) → ~21 GB (13B), calibration only")],
    """
And the implementation question, how do you compute a percentile without storing everything.
Storing every activation value to call a quantile function would need hundreds of gigabytes. Instead I keep only the top-K largest absolute values per channel, in sorted order. K is set so the buffer covers every percentile from ninety up to one hundred.
After each forward pass I concatenate the new batch and take a top-K to trim back. At the end, any percentile at or above the floor is just an index into the sorted buffer. And this is exact. It returns the identical value a full-storage quantile call would give, because the percentiles I need only ever touch the top of the distribution.
One calibration pass produces all five percentile files at once. The buffer costs about five gigabytes on the 1.3B model and twenty-one on the 13B, and only during calibration. Nothing at inference.
""")

# ---------------------------------------------------------------- save
out = os.path.join(HERE, "thesis_defense.pptx")
prs.save(out)
print("Saved:", out)
print("Slides:", len(prs.slides._sldIdLst))
